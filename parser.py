# file_parser/parser.py
import os
import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd

def parse_pdf(path):
    doc = fitz.open(path)
    text = "\n".join([page.get_text() for page in doc])
    doc.close()
    return text

def parse_docx(path):
    doc = docx.Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(path):
    prs = pptx.Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)

def parse_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def parse_file(path):
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".pdf":
        return parse_pdf(path)
    elif ext == ".docx":
        return parse_docx(path)
    elif ext == ".pptx":
        return parse_pptx(path)
    elif ext == ".csv":
        return parse_csv(path)
    elif ext in [".txt", ".md"]:
        return parse_txt(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
